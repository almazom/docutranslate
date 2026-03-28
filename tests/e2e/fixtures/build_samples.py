from __future__ import annotations

from io import BytesIO
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation


FIXTURE_DIR = Path(__file__).resolve().parent
GENERATED_DIR = FIXTURE_DIR / "generated"
SAMPLE_LINES = [
    "DocuTranslate end-to-end sample.",
    "This content should remain unchanged in skip-translate mode.",
]
RICH_DOCX_HEADING = "Quarterly brief"
RICH_DOCX_BULLETS = [
    "Budget stays lean.",
    "Formatting should survive.",
    "Translation should remain concise.",
]
RICH_XLSX_SHEETS = {
    "Summary": [
        ["Metric", "Value", "Comment"],
        ["Region", "North", "Compact office fixture"],
        ["Status", "Ready", "Short table for translation"],
    ],
    "Risks": [
        ["Risk", "Owner"],
        ["Delay", "Ops"],
        ["Budget", "PM"],
    ],
}
RICH_PPTX_TITLE = "Launch brief"
RICH_PPTX_SLIDES = [
    ("Milestones", ["Scope locked", "Pilot in April", "Public rollout in May"]),
    ("Quality gates", ["Keep layout stable", "Keep copy short"]),
]


def build_docx(path: Path) -> None:
    document = Document()
    for line in SAMPLE_LINES:
        document.add_paragraph(line)
    document.save(path)


def build_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sample"
    for row, line in enumerate(SAMPLE_LINES, start=1):
        sheet.cell(row=row, column=1, value=line)
    workbook.save(path)


def build_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = SAMPLE_LINES[0]
    slide.placeholders[1].text = SAMPLE_LINES[1]
    presentation.save(path)


def build_rich_docx(path: Path) -> None:
    document = Document()
    document.add_heading(RICH_DOCX_HEADING, level=1)
    document.add_paragraph("Compact office fixture with richer structure.")
    for bullet in RICH_DOCX_BULLETS:
        document.add_paragraph(bullet, style="List Bullet")
    table = document.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Area"
    table.rows[0].cells[1].text = "Status"
    table.rows[1].cells[0].text = "North"
    table.rows[1].cells[1].text = "Ready"
    table.rows[2].cells[0].text = "South"
    table.rows[2].cells[1].text = "Review"
    document.save(path)


def build_rich_xlsx(path: Path) -> None:
    workbook = Workbook()
    summary_sheet = workbook.active
    summary_sheet.title = "Summary"
    for row_index, row in enumerate(RICH_XLSX_SHEETS["Summary"], start=1):
        for column_index, value in enumerate(row, start=1):
            summary_sheet.cell(row=row_index, column=column_index, value=value)

    risks_sheet = workbook.create_sheet("Risks")
    for row_index, row in enumerate(RICH_XLSX_SHEETS["Risks"], start=1):
        for column_index, value in enumerate(row, start=1):
            risks_sheet.cell(row=row_index, column=column_index, value=value)
    workbook.save(path)


def build_rich_pptx(path: Path) -> None:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = RICH_PPTX_TITLE
    title_slide.placeholders[1].text = "Short multi-slide office deck."

    for slide_title, bullets in RICH_PPTX_SLIDES:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        text_frame.text = bullets[0]
        for bullet in bullets[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
    presentation.save(path)


def escape_pdf_text(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def build_pdf(path: Path) -> None:
    content_lines = ["BT", "/F1 18 Tf", "72 720 Td"]
    for index, line in enumerate(SAMPLE_LINES):
        if index:
            content_lines.append("T*")
        content_lines.append(f"({escape_pdf_text(line)}) Tj")
    content_lines.append("ET")
    stream = "\n".join(content_lines).encode("utf-8")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode("utf-8") + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    pdf = BytesIO()
    pdf.write(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]

    for index, obj in enumerate(objects, start=1):
        offsets.append(pdf.tell())
        pdf.write(f"{index} 0 obj\n".encode("utf-8"))
        pdf.write(obj)
        pdf.write(b"\nendobj\n")

    xref_offset = pdf.tell()
    pdf.write(f"xref\n0 {len(objects) + 1}\n".encode("utf-8"))
    pdf.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.write(f"{offset:010} 00000 n \n".encode("utf-8"))

    pdf.write(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("utf-8")
    )
    path.write_bytes(pdf.getvalue())


def main() -> None:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    build_docx(GENERATED_DIR / "sample.docx")
    build_rich_docx(GENERATED_DIR / "sample_rich.docx")
    build_xlsx(GENERATED_DIR / "sample.xlsx")
    build_rich_xlsx(GENERATED_DIR / "sample_rich.xlsx")
    build_pptx(GENERATED_DIR / "sample.pptx")
    build_rich_pptx(GENERATED_DIR / "sample_rich.pptx")
    build_pdf(GENERATED_DIR / "sample.pdf")


if __name__ == "__main__":
    main()
